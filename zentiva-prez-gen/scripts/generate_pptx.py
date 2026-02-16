"""
Zentiva Presentation Generator

Generates branded PowerPoint presentations using the Zentiva template.
Preserves visual elements by modifying template slides directly.
"""

import sys
import os
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from lxml import etree

# =============================================================================
# VERSION
# =============================================================================
VERSION = "v019"


def get_versioned_filename(base_name, version=None):
    """Generate a versioned filename: basename_v010.pptx"""
    v = version or VERSION
    if base_name.endswith('.pptx'):
        base_name = base_name[:-5]
    return f"{base_name}_{v}.pptx"


# =============================================================================
# ZENTIVA BRAND COLORS
# =============================================================================
ZENTIVA_DARK_BLUE = RGBColor(12, 65, 96)      # #0C4160
ZENTIVA_TEAL = RGBColor(0, 169, 143)          # #00A98F
ZENTIVA_LIGHT_TEAL = RGBColor(232, 245, 242)  # #E8F5F2

# =============================================================================
# LAYOUT NAMES
# =============================================================================
LAYOUT_TITLE = "1_Title Slide"
LAYOUT_QUOTE = "1_Title Slide"  # Same layout, has gradient shape in template slide 2
LAYOUT_CONTENT = "2_Title and Content"
LAYOUT_TWO_COLUMN = "17_Title and Content"
LAYOUT_SPLIT = "3_Title Slide"
LAYOUT_TEXT_IMAGE = "12_Title and Content"  # Title, subtitle, text left, image right
LAYOUT_HIGHLIGHT = "25_Title and Content"  # For highlight content slides with inline emphasis


def get_layout_by_name(prs, name):
    """Get a slide layout by its name."""
    for layout in prs.slide_layouts:
        if layout.name == name:
            return layout
    return prs.slide_layouts[0]


def clear_placeholder_text(slide):
    """Clear all placeholder text frames completely."""
    for shape in slide.placeholders:
        if shape.has_text_frame:
            tf = shape.text_frame
            # Clear all paragraphs
            for para in tf.paragraphs:
                para.clear()
            tf.text = ""


def delete_placeholder(slide, idx):
    """
    Delete a placeholder shape from a slide by its index.

    Args:
        slide: Slide object
        idx: Placeholder index to delete
    """
    for shape in list(slide.shapes):
        if shape.is_placeholder and shape.placeholder_format.idx == idx:
            sp = shape._element
            sp.getparent().remove(sp)
            return True
    return False


def hide_unused_placeholders(slide, keep_indices):
    """
    Hide/delete placeholders not in the keep list.

    Args:
        slide: Slide object
        keep_indices: List of placeholder indices to keep (e.g., [0, 14] for title and body)
    """
    for shape in list(slide.shapes):
        if shape.is_placeholder:
            idx = shape.placeholder_format.idx
            if idx not in keep_indices:
                # Delete the shape
                sp = shape._element
                sp.getparent().remove(sp)


def set_bullet_format(paragraph, bullet_char="•", color=None, indent_emu=None):
    """
    Set explicit bullet formatting on a paragraph.

    Args:
        paragraph: pptx paragraph object
        bullet_char: Character to use as bullet (default: •)
        color: RGBColor for bullet, or None to inherit
        indent_emu: Space between bullet and text in EMUs (default: 457200 = 0.5 inch)
    """
    pPr = paragraph._p.get_or_add_pPr()

    # Remove any buNone element that disables bullets
    buNone = pPr.find(qn('a:buNone'))
    if buNone is not None:
        pPr.remove(buNone)

    # Add bullet character
    buChar = pPr.find(qn('a:buChar'))
    if buChar is None:
        buChar = etree.SubElement(pPr, qn('a:buChar'))
    buChar.set('char', bullet_char)

    # Set bullet-to-text spacing via margin and indent
    # marL = left margin (where text starts)
    # indent = negative offset for bullet from margin
    # Space between bullet and text = marL + indent
    if indent_emu is None:
        indent_emu = 457200  # 0.5 inch in EMUs (914400 EMUs = 1 inch)

    level = paragraph.level
    base_margin = 457200 + (level * 457200)  # Base margin increases with level
    pPr.set('marL', str(base_margin))
    pPr.set('indent', str(-274320))  # Pull bullet back (creates space)

    # Set bullet color if specified
    if color:
        buClr = pPr.find(qn('a:buClr'))
        if buClr is None:
            buClr = etree.SubElement(pPr, qn('a:buClr'))
        else:
            buClr.clear()
        srgbClr = etree.SubElement(buClr, qn('a:srgbClr'))
        # RGBColor str() returns hex like "00A98F"
        srgbClr.set('val', str(color))


def fill_text_frame(shape, content, bullets=False, bullet_color=None,
                    level_styles=None):
    """
    Fill a text frame with content (string or list of bullets).

    Args:
        shape: Shape with text frame
        content: String or list of items. Items can be:
                 - str: "text" (level 0)
                 - tuple: ("text", level) where level is 0, 1, 2...
        bullets: If True, add explicit bullet formatting
        bullet_color: RGBColor for bullets, or None to use default (teal)
                     Can also be a dict {level: color} for per-level colors
        level_styles: Optional dict mapping level to style settings:
                     {0: {'font_size': Pt(20), 'font_color': BLUE, 'bullet_color': BLUE},
                      1: {'font_size': Pt(16), 'font_color': TEAL, 'bullet_color': TEAL}}
    """
    tf = shape.text_frame
    tf.text = ""

    # Default bullet color
    default_bullet_color = bullet_color if bullet_color else ZENTIVA_TEAL

    if isinstance(content, str):
        tf.paragraphs[0].text = content
    elif isinstance(content, list):
        for i, item in enumerate(content):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()

            # Parse item
            if isinstance(item, str):
                text = item
                level = 0
            elif isinstance(item, (list, tuple)) and len(item) >= 2:
                text = item[0]
                level = item[1]
            else:
                text = str(item)
                level = 0

            p.text = text
            p.level = level

            # Apply level-specific styling
            if level_styles and level in level_styles:
                style = level_styles[level]
                # Apply font styling to all runs
                for run in p.runs:
                    if 'font_size' in style:
                        run.font.size = style['font_size']
                    if 'font_color' in style:
                        run.font.color.rgb = style['font_color']
                # Get bullet color for this level
                b_color = style.get('bullet_color', default_bullet_color)
            else:
                b_color = default_bullet_color

            if bullets:
                set_bullet_format(p, "•", b_color)


def add_textbox_with_bullets(slide, left, top, width, height, content,
                              font_size=Pt(14), font_color=None, bullet_color=None):
    """
    Add a textbox with bullet-formatted content (bullets-in-textbox mode).

    Args:
        slide: Slide to add textbox to
        left, top, width, height: Position and size (Inches or Emu)
        content: List of bullet items
        font_size: Font size (default 14pt)
        font_color: Text color (default: dark blue)
        bullet_color: Bullet color (default: teal)

    Returns:
        The created textbox shape
    """
    if font_color is None:
        font_color = ZENTIVA_DARK_BLUE
    if bullet_color is None:
        bullet_color = ZENTIVA_TEAL

    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(content):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()

        if isinstance(item, str):
            p.text = item
            p.level = 0
        elif isinstance(item, (list, tuple)) and len(item) >= 2:
            p.text = item[0]
            p.level = item[1]

        # Set font formatting
        for run in p.runs:
            run.font.size = font_size
            run.font.color.rgb = font_color

        # Set explicit bullet
        set_bullet_format(p, "•", bullet_color)

    return textbox


def remove_slide(prs, index):
    """Remove a slide by index."""
    rId = prs.slides._sldIdLst[index].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[index]


# =============================================================================
# SLIDE MODIFICATION FUNCTIONS
# =============================================================================

def modify_title_slide(slide, title, subtitle=None, bg_image=None, date_str=None):
    """
    Modify an existing title slide.

    Args:
        slide: Slide object
        title: Main title text
        subtitle: Optional subtitle text
        bg_image: Optional background image path
        date_str: Optional date string (format: YYYY-MMM-DD)
    """
    clear_placeholder_text(slide)
    for shape in slide.placeholders:
        ph = shape.placeholder_format
        if ph.type == 3:  # CENTER_TITLE
            shape.text = title
        elif ph.type == 4 and subtitle:  # SUBTITLE
            shape.text = subtitle
        elif ph.type == 18 and bg_image and os.path.exists(bg_image):
            # Check if placeholder already has picture (PlaceholderPicture vs SlidePlaceholder)
            if hasattr(shape, 'insert_picture') and shape.__class__.__name__ == 'SlidePlaceholder':
                shape.insert_picture(bg_image)
            # If already a PlaceholderPicture, template image is preserved

    # Add date if provided
    if date_str:
        # Add date textbox on title slide
        # Slide is 13.33" wide x 7.5" tall
        # Position: bottom right, above the footer (footer starts ~7.0")
        date_box = slide.shapes.add_textbox(
            Inches(9.0), Inches(6.2), Inches(4.0), Inches(0.5)
        )
        tf = date_box.text_frame
        p = tf.paragraphs[0]
        p.text = date_str
        p.alignment = PP_ALIGN.RIGHT
        for run in p.runs:
            run.font.size = Pt(20)
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)  # White text


def modify_quote_slide(slide, statement):
    """Modify an existing quote/statement slide."""
    clear_placeholder_text(slide)
    for shape in slide.placeholders:
        if shape.placeholder_format.type == 3:  # CENTER_TITLE
            shape.text = statement


def modify_content_slide(slide, title, bullets):
    """Modify an existing content slide."""
    clear_placeholder_text(slide)
    if slide.shapes.title:
        slide.shapes.title.text = title
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 14:
            fill_text_frame(shape, bullets)
            break


def modify_two_column_slide(slide, title, left_content, right_content):
    """Modify an existing two-column slide."""
    clear_placeholder_text(slide)
    if slide.shapes.title:
        slide.shapes.title.text = title
    filled = set()
    for shape in slide.placeholders:
        ph = shape.placeholder_format
        if ph.idx == 14 and 14 not in filled:
            fill_text_frame(shape, left_content)
            filled.add(14)
        elif ph.idx == 15 and 15 not in filled:
            fill_text_frame(shape, right_content)
            filled.add(15)


def modify_split_slide(slide, title=None, subtitle=None, image_path=None):
    """Modify an existing split layout slide."""
    clear_placeholder_text(slide)
    for shape in slide.placeholders:
        ph = shape.placeholder_format
        if ph.type == 3 and title:
            shape.text = title
        elif ph.type == 4 and subtitle:
            shape.text = subtitle
        elif ph.type == 18 and image_path and os.path.exists(image_path):
            # Check if placeholder already has picture
            if hasattr(shape, 'insert_picture') and shape.__class__.__name__ == 'SlidePlaceholder':
                shape.insert_picture(image_path)
            # If already a PlaceholderPicture, template image is preserved


# =============================================================================
# ADD NEW SLIDES (for slides not in template)
# =============================================================================

def add_content_slide(prs, title, bullets, subtitle=None, mode='placeholder'):
    """
    Add a new content slide with bullet points.

    Args:
        prs: Presentation object
        title: Slide title
        bullets: List of bullet items
        subtitle: Optional green subtitle below title
        mode: 'placeholder' = use body placeholder for bullets (default)
              'textbox' = use free-form textbox for bullets

    Returns:
        The created slide
    """
    layout = get_layout_by_name(prs, LAYOUT_CONTENT)
    slide = prs.slides.add_slide(layout)

    # Set title
    if slide.shapes.title:
        slide.shapes.title.text = title

    # Define level-based styling:
    # Level 0 (parent): bigger font, dark blue text + bullets
    # Level 1+ (child): smaller font, dark blue text, teal/green bullets
    level_styles = {
        0: {
            'font_size': Pt(20),
            'font_color': ZENTIVA_DARK_BLUE,
            'bullet_color': ZENTIVA_DARK_BLUE
        },
        1: {
            'font_size': Pt(16),
            'font_color': ZENTIVA_TEAL,  # Green text
            'bullet_color': ZENTIVA_TEAL  # Green bullet
        },
        2: {
            'font_size': Pt(14),
            'font_color': ZENTIVA_TEAL,  # Green text
            'bullet_color': ZENTIVA_TEAL  # Green bullet
        }
    }

    if mode == 'placeholder':
        # Keep only title (idx=0) and body (idx=14), remove others like OBJECT placeholder
        hide_unused_placeholders(slide, keep_indices=[0, 14])

        # Fill the body placeholder with subtitle (if any) + bullets
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 14:
                fill_text_frame_with_subtitle(
                    shape, bullets, subtitle=subtitle, level_styles=level_styles
                )
                break

    elif mode == 'textbox':
        # Remove all content placeholders, keep only title
        hide_unused_placeholders(slide, keep_indices=[0])

        # Add textbox with subtitle + bullets using level styles
        add_textbox_with_subtitle_and_bullets(
            slide,
            left=Inches(0.5),
            top=Inches(1.5),
            width=Inches(9),
            height=Inches(5),
            subtitle=subtitle,
            content=bullets,
            level_styles=level_styles
        )

    return slide


def fill_text_frame_with_subtitle(shape, content, subtitle=None, level_styles=None):
    """
    Fill a text frame with optional subtitle and bulleted content.

    Args:
        shape: Shape with text frame
        content: List of bullet items
        subtitle: Optional subtitle (green, no bullet)
        level_styles: Dict mapping level to style settings
    """
    tf = shape.text_frame
    tf.text = ""

    para_index = 0

    # Add subtitle first (no bullet, green/teal color)
    if subtitle:
        p = tf.paragraphs[0]
        p.text = subtitle
        p.level = 0
        # Style subtitle: green, slightly smaller than parent bullets
        for run in p.runs:
            run.font.size = Pt(18)
            run.font.color.rgb = ZENTIVA_TEAL
        # No bullet for subtitle
        pPr = p._p.get_or_add_pPr()
        buNone = etree.SubElement(pPr, qn('a:buNone'))
        # Add space after subtitle (18pt = ~0.25 inch)
        p.space_after = Pt(18)
        para_index = 1

    # Add bullet content
    for i, item in enumerate(content):
        if para_index == 0 and i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Parse item
        if isinstance(item, str):
            text = item
            level = 0
        elif isinstance(item, (list, tuple)) and len(item) >= 2:
            text = item[0]
            level = item[1]
        else:
            text = str(item)
            level = 0

        p.text = text
        p.level = level

        # Apply level-specific styling
        if level_styles and level in level_styles:
            style = level_styles[level]
            for run in p.runs:
                if 'font_size' in style:
                    run.font.size = style['font_size']
                if 'font_color' in style:
                    run.font.color.rgb = style['font_color']
            b_color = style.get('bullet_color', ZENTIVA_TEAL)
        else:
            b_color = ZENTIVA_TEAL

        set_bullet_format(p, "•", b_color)


def add_textbox_with_subtitle_and_bullets(slide, left, top, width, height, content,
                                           subtitle=None, level_styles=None):
    """
    Add a textbox with optional subtitle and level-styled bullet content.

    Args:
        slide: Slide to add textbox to
        left, top, width, height: Position and size
        content: List of bullet items
        subtitle: Optional subtitle (green, no bullet)
        level_styles: Dict mapping level to style settings
    """
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True

    para_index = 0

    # Add subtitle first (no bullet, green/teal color)
    if subtitle:
        p = tf.paragraphs[0]
        p.text = subtitle
        p.level = 0
        for run in p.runs:
            run.font.size = Pt(18)
            run.font.color.rgb = ZENTIVA_TEAL
        # No bullet for subtitle
        pPr = p._p.get_or_add_pPr()
        buNone = etree.SubElement(pPr, qn('a:buNone'))
        para_index = 1

    # Add bullet content
    for i, item in enumerate(content):
        if para_index == 0 and i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Parse item
        if isinstance(item, str):
            text = item
            level = 0
        elif isinstance(item, (list, tuple)) and len(item) >= 2:
            text = item[0]
            level = item[1]
        else:
            text = str(item)
            level = 0

        p.text = text
        p.level = level

        # Get style for this level
        if level_styles:
            style = level_styles.get(level, level_styles.get(0, {}))
        else:
            style = {}

        # Apply font formatting
        for run in p.runs:
            if 'font_size' in style:
                run.font.size = style['font_size']
            if 'font_color' in style:
                run.font.color.rgb = style['font_color']

        # Set bullet
        b_color = style.get('bullet_color', ZENTIVA_TEAL)
        set_bullet_format(p, "•", b_color)

    return textbox


def add_two_column_slide(prs, title, left_content, right_content, mode='placeholder'):
    """
    Add a new two-column slide.

    Args:
        prs: Presentation object
        title: Slide title
        left_content: Left column bullet items
        right_content: Right column bullet items
        mode: 'placeholder' = use body placeholders for bullets (default)
              'textbox' = use free-form textboxes for bullets

    Returns:
        The created slide
    """
    layout = get_layout_by_name(prs, LAYOUT_TWO_COLUMN)
    slide = prs.slides.add_slide(layout)

    # Set title
    if slide.shapes.title:
        slide.shapes.title.text = title

    if mode == 'placeholder':
        # Keep only title (idx=0), left body (idx=14), right body (idx=15)
        hide_unused_placeholders(slide, keep_indices=[0, 14, 15])

        # Fill the body placeholders
        for shape in slide.placeholders:
            ph = shape.placeholder_format
            if ph.idx == 14:
                # Left column: teal bullets on white background
                fill_text_frame(shape, left_content, bullets=True, bullet_color=ZENTIVA_TEAL)
            elif ph.idx == 15:
                # Right column: white bullets on dark gradient circle
                fill_text_frame(shape, right_content, bullets=True,
                              bullet_color=RGBColor(255, 255, 255))

    elif mode == 'textbox':
        # Remove all content placeholders, keep only title
        hide_unused_placeholders(slide, keep_indices=[0])

        # Add left textbox
        add_textbox_with_bullets(
            slide,
            left=Inches(0.5),
            top=Inches(1.5),
            width=Inches(4.5),
            height=Inches(5),
            content=left_content,
            font_size=Pt(16),
            font_color=ZENTIVA_DARK_BLUE,
            bullet_color=ZENTIVA_TEAL
        )

        # Add right textbox
        add_textbox_with_bullets(
            slide,
            left=Inches(5.2),
            top=Inches(1.5),
            width=Inches(4.5),
            height=Inches(5),
            content=right_content,
            font_size=Pt(16),
            font_color=RGBColor(255, 255, 255),
            bullet_color=RGBColor(255, 255, 255)
        )

    return slide


def add_text_image_slide(prs, title, subtitle, content, image_path=None):
    """
    Add a slide with title, subtitle, text content on left, and image on right.
    Based on template slide 12 / layout "12_Title and Content".

    Args:
        prs: Presentation object
        title: Main title (blue)
        subtitle: Subtitle text (green/teal)
        content: List of bullet items for left side
        image_path: Path to image for right placeholder (optional)

    Returns:
        The created slide
    """
    layout = get_layout_by_name(prs, LAYOUT_TEXT_IMAGE)
    slide = prs.slides.add_slide(layout)

    # Keep only needed placeholders: title(0), subtitle(14), content(13), picture(15)
    hide_unused_placeholders(slide, keep_indices=[0, 13, 14, 15])

    # Set title (blue)
    if slide.shapes.title:
        slide.shapes.title.text = title

    # Define level styles for content bullets
    level_styles = {
        0: {
            'font_size': Pt(18),
            'font_color': ZENTIVA_DARK_BLUE,
            'bullet_color': ZENTIVA_DARK_BLUE
        },
        1: {
            'font_size': Pt(14),
            'font_color': ZENTIVA_TEAL,  # Green text
            'bullet_color': ZENTIVA_TEAL  # Green bullet
        }
    }

    for shape in slide.placeholders:
        ph = shape.placeholder_format

        # Subtitle placeholder (idx=14) - green text, no bullet
        if ph.idx == 14:
            shape.text_frame.text = ""
            p = shape.text_frame.paragraphs[0]
            p.text = subtitle
            for run in p.runs:
                run.font.size = Pt(16)
                run.font.color.rgb = ZENTIVA_TEAL

        # Content placeholder (idx=13) - bullets on left
        elif ph.idx == 13:
            fill_text_frame(shape, content, bullets=True, level_styles=level_styles)

        # Picture placeholder (idx=15)
        elif ph.idx == 15 and image_path and os.path.exists(image_path):
            # Check if it's a SlidePlaceholder that can accept an image
            if shape.__class__.__name__ == 'SlidePlaceholder':
                shape.insert_picture(image_path)

    return slide


def parse_highlight_text(text):
    """
    Parse text with <<highlight>> markers into segments.

    Args:
        text: String with <<highlighted>> portions marked

    Returns:
        List of (text, is_highlighted) tuples

    Example:
        "Normal <<highlighted>> normal" ->
        [("Normal ", False), ("highlighted", True), (" normal", False)]

        "Value is <<>50 million>> here" ->
        [("Value is ", False), (">50 million", True), (" here", False)]
    """
    import re
    segments = []
    # Use non-greedy match to handle > inside highlighted text
    # Match << followed by any chars (non-greedy) followed by >>
    pattern = r'<<(.+?)>>'

    last_end = 0
    for match in re.finditer(pattern, text):
        # Add normal text before the match
        if match.start() > last_end:
            segments.append((text[last_end:match.start()], False))
        # Add highlighted text
        segments.append((match.group(1), True))
        last_end = match.end()

    # Add remaining normal text
    if last_end < len(text):
        segments.append((text[last_end:], False))

    # If no highlights found, return whole text as normal
    if not segments:
        segments = [(text, False)]

    return segments


def add_highlight_slide(prs, title, content):
    """
    Add a slide with inline text highlighting (teal emphasis on key phrases).

    This replicates the format from slide 16 of Brand_v002.pptx where specific
    text segments within bullet points are highlighted in teal color.

    Args:
        prs: Presentation object
        title: Slide title (dark blue)
        content: List of bullet items. Each item can be:
            - str: "Normal text with <<highlighted>> portions" (level 0)
            - tuple: ("Text with <<highlights>>", level) where level is 0, 1, 2...

        Use <<text>> syntax to mark portions that should be highlighted in teal.

    Returns:
        The created slide

    Example:
        content = [
            "<<Key point>> with additional context",
            ("Supporting detail with <<emphasis>>", 1),
            "Another point <<highlighted phrase>> continues here"
        ]
    """
    # Try to get the highlight layout, fall back to content layout
    layout = get_layout_by_name(prs, LAYOUT_HIGHLIGHT)
    if layout is None or layout.name != LAYOUT_HIGHLIGHT:
        layout = get_layout_by_name(prs, LAYOUT_CONTENT)

    slide = prs.slides.add_slide(layout)

    # Keep only title placeholder, remove others
    hide_unused_placeholders(slide, keep_indices=[0])

    # Set title
    if slide.shapes.title:
        slide.shapes.title.text = title

    # Create textbox for content (not using placeholder for better control)
    # Position similar to template: below title with margins
    textbox = slide.shapes.add_textbox(
        Inches(0.6), Inches(1.4), Inches(12.0), Inches(5.5)
    )
    tf = textbox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(content):
        # Parse item
        if isinstance(item, str):
            text = item
            level = 0
        elif isinstance(item, (list, tuple)) and len(item) >= 2:
            text = item[0]
            level = item[1]
        else:
            text = str(item)
            level = 0

        # Get or create paragraph
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.level = level

        # Parse highlight segments
        segments = parse_highlight_text(text)

        # Determine styling based on level
        is_level_0 = level == 0
        base_font_size = Pt(18) if is_level_0 else Pt(16)
        is_bold = is_level_0

        # Clear any default text and add runs for each segment
        p.clear()

        for seg_text, is_highlighted in segments:
            run = p.add_run()
            run.text = seg_text
            run.font.size = base_font_size
            run.font.bold = is_bold

            if is_highlighted:
                run.font.color.rgb = ZENTIVA_TEAL  # Highlighted = teal
            else:
                run.font.color.rgb = ZENTIVA_DARK_BLUE  # Normal = dark blue

        # Set bullet formatting
        bullet_color = ZENTIVA_DARK_BLUE if is_level_0 else ZENTIVA_TEAL
        set_bullet_format(p, "•", bullet_color)

    return slide


def move_slide_to_end(prs, slide_index):
    """
    Move a slide to the end of the presentation.

    Args:
        prs: Presentation object
        slide_index: Index of slide to move (0-based)
    """
    # Get the slide element
    slide_id = prs.slides._sldIdLst[slide_index]
    # Remove from current position
    prs.slides._sldIdLst.remove(slide_id)
    # Append to end
    prs.slides._sldIdLst.append(slide_id)


def hide_slide(prs, slide_index):
    """
    Hide a slide from the presentation (set show attribute to false).

    Args:
        prs: Presentation object
        slide_index: Index of slide to hide (0-based)
    """
    slide = prs.slides[slide_index]
    # Access the slide element and set show="0"
    slide._element.set('show', '0')


def add_slide_number(slide, current, total):
    """
    Add a slide number to a slide in "current / total" format.

    Args:
        slide: Slide object
        current: Current slide number
        total: Total number of slides
    """
    number_text = f"{current} / {total}"

    # Check if slide number placeholder exists and fill it
    for shape in slide.placeholders:
        if shape.placeholder_format.type == 13:  # SLIDE_NUMBER
            shape.text = number_text
            for para in shape.text_frame.paragraphs:
                para.alignment = PP_ALIGN.RIGHT
                for run in para.runs:
                    run.font.color.rgb = ZENTIVA_DARK_BLUE
            return

    # If no placeholder, add textbox at far bottom right
    # Slide is 13.33" wide, position near right edge with small margin
    num_box = slide.shapes.add_textbox(
        Inches(11.8), Inches(6.9), Inches(1.3), Inches(0.3)
    )
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = number_text
    p.alignment = PP_ALIGN.RIGHT
    for run in p.runs:
        run.font.size = Pt(12)
        run.font.color.rgb = ZENTIVA_DARK_BLUE


def add_slide_numbers_to_presentation(prs, start_from_slide=1, total_visible=None):
    """
    Add slide numbers to all slides except the first one (title slide).

    Args:
        prs: Presentation object
        start_from_slide: Index of first slide to number (default 1, skips title)
        total_visible: Total visible slides count (excludes hidden). If None, calculated.
    """
    if total_visible is None:
        total_visible = len(prs.slides) - start_from_slide

    for i, slide in enumerate(prs.slides):
        if i >= start_from_slide:
            current = i - start_from_slide + 1
            add_slide_number(slide, current, total_visible)


def add_conclusion_slide(prs, title="Conclusions", subtitle=None, takeaways=None):
    """
    Add a conclusion/summary slide.

    Args:
        prs: Presentation object
        title: Slide title (default: "Key Takeaways")
        subtitle: Optional subtitle
        takeaways: List of key takeaway points

    Returns:
        The created slide
    """
    if takeaways is None:
        takeaways = [
            "Summary point one",
            "Summary point two",
            "Summary point three"
        ]

    # Use the content layout
    layout = get_layout_by_name(prs, LAYOUT_CONTENT)
    slide = prs.slides.add_slide(layout)

    # Keep only needed placeholders
    hide_unused_placeholders(slide, keep_indices=[0, 14])

    # Set title
    if slide.shapes.title:
        slide.shapes.title.text = title

    # Level styles for conclusion: all at level 0 with larger font
    level_styles = {
        0: {
            'font_size': Pt(22),
            'font_color': ZENTIVA_DARK_BLUE,
            'bullet_color': ZENTIVA_TEAL
        }
    }

    # Fill content with subtitle + takeaways
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 14:
            fill_text_frame_with_subtitle(
                shape, takeaways, subtitle=subtitle, level_styles=level_styles
            )
            break

    return slide


def get_current_date_formatted():
    """Return current date in YYYY-MMM-DD format (e.g., 2026-Feb-14)."""
    return datetime.now().strftime("%Y-%b-%d")


# =============================================================================
# MAIN GENERATION
# =============================================================================

def generate_presentation(output_path, template_path, slides_spec):
    """
    Generate a presentation by modifying template slides.

    Args:
        output_path: Path to save the presentation
        template_path: Path to Brand_v001.pptx template
        slides_spec: List of slide specifications, each with:
            - 'type': 'title', 'quote', 'content', 'two_column', 'split'
            - 'title': Slide title/statement
            - 'subtitle': Optional subtitle
            - 'content': Bullets list
            - 'content2': Second column content
            - 'image': Image path
    """
    prs = Presentation(template_path)

    # Template has 14 slides. We'll modify the ones we need and remove the rest.
    # Slide indices (0-based):
    # 0: Title slide (1_Title Slide) - has dark bg
    # 1: Quote slide (1_Title Slide) - has gradient shape
    # 2-5: Split layouts
    # 6+: Various content layouts

    # Process each slide spec
    slides_to_keep = []

    for i, spec in enumerate(slides_spec):
        slide_type = spec.get('type', 'content')

        if slide_type == 'title' and i == 0:
            # Modify template slide 0
            modify_title_slide(prs.slides[0], spec.get('title', ''),
                             spec.get('subtitle'), spec.get('image'))
            slides_to_keep.append(0)

        elif slide_type == 'quote':
            # Modify template slide 1 (has gradient)
            if 1 not in slides_to_keep:
                modify_quote_slide(prs.slides[1], spec.get('title', ''))
                slides_to_keep.append(1)
            else:
                # Need to add additional quote slide
                add_content_slide(prs, spec.get('title', ''), [])

        elif slide_type == 'content':
            add_content_slide(prs, spec.get('title', ''), spec.get('content', []))

        elif slide_type == 'two_column':
            add_two_column_slide(prs, spec.get('title', ''),
                               spec.get('content', []), spec.get('content2', []))

        elif slide_type == 'highlight':
            add_highlight_slide(prs, spec.get('title', ''), spec.get('content', []))

        elif slide_type == 'split':
            # Modify template slide 2 (split layout)
            if 2 not in slides_to_keep:
                modify_split_slide(prs.slides[2], spec.get('title'),
                                 spec.get('subtitle'), spec.get('image'))
                slides_to_keep.append(2)

    # Remove unused template slides (in reverse order to maintain indices)
    for i in range(len(prs.slides) - 1, -1, -1):
        # Keep slides we modified and any slides we added (which are at end)
        if i not in slides_to_keep and i < 14:  # 14 is original template count
            remove_slide(prs, i)

    prs.save(output_path)
    print(f"[zentiva-prez-gen {VERSION}] Presentation saved to {output_path}")
    return output_path


def generate_test_presentation(output_path, template_path):
    """Generate a test presentation showcasing all layout types."""
    prs = Presentation(template_path)

    script_dir = os.path.dirname(os.path.abspath(__file__))
    assets_dir = os.path.join(script_dir, "..", "assets")
    title_bg = os.path.join(assets_dir, "title_bg.png")

    # Modify slide 0 (Title) with date
    modify_title_slide(prs.slides[0], "Test Presentation",
                       "Verification of All Layouts",
                       title_bg if os.path.exists(title_bg) else None,
                       date_str=get_current_date_formatted())

    # Modify slide 1 (Quote with gradient) - will be moved to end and hidden
    modify_quote_slide(prs.slides[1],
                       "This is a test quote to verify the gradient background")

    # Modify slide 2 (Split layout) - will be moved to end and hidden
    modify_split_slide(prs.slides[2], "Split Layout Test", "With Right Image",
                       title_bg if os.path.exists(title_bg) else None)

    # Add content slide with subtitle and hierarchical styling
    content_slide = add_content_slide(prs, "Content Slide Test",
        bullets=[
            "First level bullet point",
            ("Second level indented", 1),
            ("Another second level", 1),
            "Back to first level",
            ("Nested point", 1)
        ],
        subtitle="This is a green subtitle for the content slide"
    )

    # Add two-column slide
    two_col_slide = add_two_column_slide(prs, "Two Column Layout Test",
        ["Left column point 1", "Left column point 2", "Left column point 3"],
        ["Right column point 1", "Right column point 2", "Right column point 3"])

    # Add text+image slide (layout 12)
    text_image_slide = add_text_image_slide(prs,
        title="Text and Image Layout",
        subtitle="Slide 12 style with content left and image right",
        content=[
            "Key feature one",
            ("Detail about feature one", 1),
            "Key feature two",
            ("Detail about feature two", 1),
            "Key feature three"
        ],
        image_path=title_bg if os.path.exists(title_bg) else None
    )

    # Add highlight slide (new in v019)
    highlight_slide = add_highlight_slide(prs,
        title="Highlight Slide Test",
        content=[
            "<<Key benefit>> with supporting explanation text",
            ("Detail point with <<inline emphasis>> in the middle", 1),
            "Another <<important point>> that drives the message",
            ("CPRD contains <<>50 million patients>> and covers ~24%", 1),
            "<<Multiple highlights>> can appear <<in one line>>",
        ]
    )

    # Add conclusion slide
    conclusion_slide = add_conclusion_slide(prs,
        title="Conclusions",
        subtitle="Summary of this presentation",
        takeaways=[
            "Zentiva branding is fully preserved",
            "Multiple slide layouts are supported",
            "Hierarchical bullet points with color styling",
            "Text and image layouts available",
            "Highlight slides with inline emphasis (v019)",
            "Slide management (hide, reorder) implemented"
        ]
    )

    # Remove unused template slides (slides 3-13)
    # After this: 0=Title, 1=Quote, 2=Split, 3+=added slides
    for i in range(13, 2, -1):
        remove_slide(prs, i)

    # Now reorder: Move Quote (1) and Split (2) to the end and hide them
    # After removals: 0=Title, 1=Quote, 2=Split, 3=Content, 4=TwoCol, 5=TextImage, 6=Highlight, 7=Conclusion
    # Move Quote (index 1) to end
    move_slide_to_end(prs, 1)
    # Now: 0=Title, 1=Split, 2=Content, 3=TwoCol, 4=TextImage, 5=Highlight, 6=Conclusion, 7=Quote
    # Move Split (now index 1) to end
    move_slide_to_end(prs, 1)
    # Now: 0=Title, 1=Content, 2=TwoCol, 3=TextImage, 4=Highlight, 5=Conclusion, 6=Quote, 7=Split

    # Hide the last two slides (Quote and Split)
    hide_slide(prs, len(prs.slides) - 2)  # Quote
    hide_slide(prs, len(prs.slides) - 1)  # Split

    # Add slide numbers to all visible slides except title (index 0)
    # Visible slides: 0=Title, 1=Content, 2=TwoCol, 3=TextImage, 4=Highlight, 5=Conclusion
    # Hidden slides at end: 6=Quote, 7=Split
    visible_count = len(prs.slides) - 2  # Exclude 2 hidden slides
    total_numbered = visible_count - 1   # Exclude title slide from numbering
    for i in range(1, visible_count):    # Skip title (0), stop before hidden
        current = i  # Slide 1 = "1", Slide 2 = "2", etc.
        add_slide_number(prs.slides[i], current, total_numbered)

    prs.save(output_path)
    print(f"[zentiva-prez-gen {VERSION}] Test presentation saved to {output_path}")
    return output_path


if __name__ == "__main__":
    script_dir = os.path.dirname(os.path.abspath(__file__))
    assets_dir = os.path.join(script_dir, "..", "assets")
    template = os.path.join(assets_dir, "Brand_v001.pptx")

    if len(sys.argv) < 2:
        print(f"Zentiva Presentation Generator {VERSION}")
        print("Usage: python generate_pptx.py --test [output_name]")
        print("       python generate_pptx.py --version")
        sys.exit(1)

    if sys.argv[1] == "--version":
        print(f"zentiva-prez-gen {VERSION}")
        sys.exit(0)

    if sys.argv[1] == "--test":
        base_name = sys.argv[2] if len(sys.argv) > 2 else "test_output"
        output_path = get_versioned_filename(base_name)
        generate_test_presentation(output_path, template)
